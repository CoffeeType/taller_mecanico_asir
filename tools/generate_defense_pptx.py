#!/usr/bin/env python3
"""
Genera la presentación de defensa del PFC (castellano, ~20 min) a partir del COPIAPFC.
Requisitos: python-pptx, pillow  →  python -m pip install python-pptx pillow

Uso (desde la raíz del repo):
    python tools/generate_defense_pptx.py
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Corporate Blue (ppt-visual)
COLOR_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
COLOR_SECONDARY = RGBColor(0x34, 0x98, 0xDB)
COLOR_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)

FONT_TITLE = "Montserrat"
FONT_BODY = "Open Sans"
FONT_FALLBACK = "Calibri"

COPIAPFC_NAME = "COPIAPFC_Taller_Mecanico_ASIR_Antonio_Corredera_Cubells_con_diagramas.docx"
OUTPUT_NAME = "PFC_Defensa_Taller_Mecanico_ASIR_Antonio_Corredera_Cubells.pptx"


def repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def resolve_image(root: Path, filename: str) -> Path | None:
    """Busca PNG en carpeta descomprimida del docx o extrae del .docx."""
    dirs = [
        root / "docs" / "_copiapfc_work" / "word" / "media",
        root / "docs" / "_reaudit" / "word" / "media",
        root / "docs" / "_copiapfc_unzip" / "word" / "media",
    ]
    for d in dirs:
        p = d / filename
        if p.is_file():
            return p
    docx = root / "docs" / COPIAPFC_NAME
    if not docx.is_file():
        return None
    inner = f"word/media/{filename}"
    cache = root / "docs" / ".pptx_gen_media"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / filename
    with zipfile.ZipFile(docx) as z:
        if inner not in z.namelist():
            return None
        out.write_bytes(z.read(inner))
    return out


def _set_font_safe(run, name: str, size_pt: int, bold: bool = False, color: RGBColor | None = None):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    try:
        run._element.rPr.rFonts.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ascii", name)
        run._element.rPr.rFonts.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}hAnsi", name)
    except Exception:
        pass


def style_title_shape(shape, text: str):
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            _set_font_safe(r, FONT_TITLE, 28, bold=True, color=COLOR_PRIMARY)
    if not tf.paragraphs[0].runs:
        run = tf.paragraphs[0].add_run()
        _set_font_safe(run, FONT_TITLE, 28, bold=True, color=COLOR_PRIMARY)


def style_body_frame(tf, bullets: list[str]):
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(6)
        p.font.size = Pt(16)
        p.font.name = FONT_BODY
        p.font.color.rgb = COLOR_TEXT
        for r in p.runs:
            _set_font_safe(r, FONT_BODY, 16, color=COLOR_TEXT)


def set_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text
    for p in notes.paragraphs:
        p.font.size = Pt(11)
        p.font.name = FONT_FALLBACK
        p.font.color.rgb = COLOR_TEXT


def fill_title_slide(slide, title: str, subtitle_lines: list[str]):
    slide.shapes.title.text = title
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            _set_font_safe(r, FONT_TITLE, 36, bold=True, color=COLOR_PRIMARY)
    sub = slide.placeholders[1]
    sub.text = "\n".join(subtitle_lines)
    stf = sub.text_frame
    for p in stf.paragraphs:
        for r in p.runs:
            _set_font_safe(r, FONT_BODY, 18, color=COLOR_TEXT)


def add_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    style_title_shape(slide.shapes.title, title)
    body = slide.placeholders[1]
    style_body_frame(body.text_frame, bullets)
    set_notes(slide, notes)
    return slide


def add_picture_slide(
    prs: Presentation,
    title: str,
    image_path: Path | None,
    bullets: list[str],
    notes: str,
    left_in: float = 0.45,
    top_in: float = 1.25,
    width_in: float = 12.0,
):
    layout = prs.slide_layouts[5]  # Title Only — luego añadimos cuerpo manual
    slide = prs.slides.add_slide(layout)
    style_title_shape(slide.shapes.title, title)
    title_bottom = Inches(1.15)
    if image_path and image_path.is_file():
        slide.shapes.add_picture(str(image_path), Inches(left_in), Inches(top_in), width=Inches(width_in))
    # Bullets a la derecha o debajo si hay imagen grande: debajo
    box_left = Inches(0.5)
    box_top = Inches(5.85)
    box_w = Inches(12.3)
    box_h = Inches(1.35)
    shape = slide.shapes.add_textbox(box_left, box_top, box_w, box_h)
    style_body_frame(shape.text_frame, bullets)
    set_notes(slide, notes)
    return slide


def add_two_picture_slide(
    prs: Presentation,
    title: str,
    path_left: Path | None,
    path_right: Path | None,
    notes: str,
):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    style_title_shape(slide.shapes.title, title)
    w = Inches(5.9)
    top = Inches(1.35)
    if path_left and path_left.is_file():
        slide.shapes.add_picture(str(path_left), Inches(0.45), top, width=w)
    if path_right and path_right.is_file():
        slide.shapes.add_picture(str(path_right), Inches(6.85), top, width=w)
    set_notes(slide, notes)
    return slide


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], notes: str):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    style_title_shape(slide.shapes.title, title)
    ncols = len(headers)
    nrows = 1 + len(rows)
    left, top = Inches(0.5), Inches(1.35)
    width, height = Inches(12.3), Inches(0.35 * nrows + 0.2)
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = FONT_FALLBACK
            p.font.color.rgb = COLOR_PRIMARY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = FONT_FALLBACK
                p.font.color.rgb = COLOR_TEXT
    set_notes(slide, notes)
    return slide


def set_slide_background_light(slide):
    """Fondo gris muy claro en la parte del master (aprox vía fill del slide)."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def build_presentation(root: Path) -> Presentation:
    prs = Presentation()
    prs.core_properties.title = "Defensa PFC — Taller mecánico (ASIR)"
    prs.core_properties.author = "Antonio Corredera Cubells"

    img1 = resolve_image(root, "image1.png")
    img2 = resolve_image(root, "image2.png")
    img3 = resolve_image(root, "image3.png")
    img4 = resolve_image(root, "image4.png")
    img19 = resolve_image(root, "image19.png")
    img20 = resolve_image(root, "image20.png")

    # 1 — Portada
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background_light(slide0)
    fill_title_slide(
        slide0,
        "Implantación y monitorización de un sistema web\npara la gestión de un taller mecánico",
        [
            "Proyecto final de ciclo — CFGS Administración de Sistemas Informáticos en Red",
            "Autor: Antonio Corredera Cubells",
            "Curso académico: 2025/2026",
            "Centro docente: [indicar centro]",
            "Tutor/a individual: [indicar tutor/a]",
        ],
    )
    set_notes(
        slide0,
        "0:00–0:30. Saludo, nombre y título. Una frase: el trabajo integra aplicación web LAMP, "
        "contenedores, despliegue en AWS y observabilidad con Prometheus y Grafana. No leer la portada.",
    )

    # 2 — Agenda
    add_content_slide(
        prs,
        "Estructura de la exposición (~20 min)",
        [
            "Bloque 1 — Problema, objetivos y encaje en el ciclo ASIR (≈5 min)",
            "Bloque 2 — Diseño, aplicación, datos, Docker, seguridad y AWS (≈7 min)",
            "Bloque 3 — Monitorización, evidencias y validación (≈6 min)",
            "Cierre — Presupuesto, conclusiones, mejoras y turno de preguntas (≈2 min)",
        ],
        "0:30–1:00. Anuncia los tres bloques del documento (7.4). Indica que habrá demo en navegador, "
        "terminal y Grafana. Reserva tiempo al final para preguntas.",
    )

    # 3 — Problema
    add_content_slide(
        prs,
        "Problema y alcance del proyecto",
        [
            "Un taller mecánico necesita ordenar citas, comunicación con clientes y administración interna.",
            "Se implantó una solución web con perfiles visitante, usuario registrado y administrador.",
            "El foco ASIR va más allá del CRUD: infraestructura reproducible, seguridad, cloud y métricas.",
        ],
        "1:00–1:30. Contexto breve: no es solo 'una web', sino un servicio desplegable y observable. "
        "Conecta con la memoria (introducción).",
    )

    # 4 — Objetivos
    add_content_slide(
        prs,
        "Objetivos del PFC",
        [
            "Diseñar e implantar una aplicación web funcional para la gestión básica de un taller mecánico.",
            "Definir una base de datos relacional (usuarios, citas, noticias, consejos) con integridad referencial.",
            "Publicar la solución con Docker para despliegue reproducible en local y en cloud.",
            "Añadir monitorización: salud del sistema, uso de la aplicación, BD y métricas de negocio.",
            "Documentar despliegue profesional en AWS con costes, seguridad perimetral y mantenimiento.",
        ],
        "1:30–2:00. Lista corta; no leer palabra por palabra. Subraya integración ciclo completo.",
    )

    # 5 — Módulos ASIR (tabla)
    add_table_slide(
        prs,
        "Competencias del CFGS ASIR aplicadas",
        ["Módulo", "Contenidos aplicados en el PFC"],
        [
            [
                "Implantación de Aplicaciones Web",
                "PHP, Apache, sesiones, MVC ligero, formularios, panel admin, API citas.",
            ],
            [
                "Planificación y Administración de Redes",
                "Docker Compose, redes, volúmenes, DNS interno, reglas de entrada en AWS (SG).",
            ],
            [
                "Proyecto de administración de sistemas",
                "Planificación, documentación, despliegue EC2, observabilidad y pruebas.",
            ],
        ],
        "2:00–2:45. Relaciona cada módulo con lo que enseñarás después. Coherencia con Tabla 1 de la memoria.",
    )

    # 6 — Arquitectura
    add_picture_slide(
        prs,
        "Arquitectura general (Fig. 1)",
        img1,
        [
            "Capa web (Apache/PHP) y persistencia MySQL.",
            "Stack Prometheus / Grafana / Alertmanager y exportadores.",
            "Simulador de carga opcional para evidencia en defensa.",
        ],
        "2:45–3:30. Señala flujo de datos y fronteras. Menciona que el diagrama está en la memoria.",
        width_in=11.2,
        left_in=0.9,
        top_in=1.2,
    )

    # 7 — Perfiles + demo
    add_table_slide(
        prs,
        "Aplicación web — perfiles y funciones principales",
        ["Perfil", "Funciones (resumen)"],
        [
            ["Visitante", "Inicio, noticias, registro, login, solicitud de cita (invitado)."],
            ["Usuario", "Gestión de perfil y citas en entorno autenticado."],
            ["Administrador", "Gestión de citas, noticias, usuarios y utilidades de operación."],
        ],
        "3:30–5:00. DEMO: navegador — flujo corto login + cita + vista admin si aplica. "
        "Habla sobre valor para el taller, no sobre cada fichero PHP.",
    )

    # 8 — Datos y flujo
    add_two_picture_slide(
        prs,
        "Modelo de datos y flujo de reserva (Fig. 2 y 3)",
        img2,
        img3,
        "5:00–5:45. Explica entidades clave y flujo usuario→cita. No entrar en cada columna SQL.",
    )

    # 9 — Docker
    add_picture_slide(
        prs,
        "Implantación con Docker Compose",
        img20,
        [
            "Servicios: web, MySQL, Prometheus, Grafana, Alertmanager, exportadores según compose.",
            "Healthchecks: el servicio web espera a BD lista; evita fallos en arranque.",
            "MySQL sin publicar puertos al host en producción; reduce superficie de ataque.",
        ],
        "5:45–6:45. DEMO: docker compose ps. Explica por qué no exponer MySQL. Enlaza con competencias ASIR.",
        left_in=0.5,
        top_in=1.25,
        width_in=7.2,
    )

    # 10 — Seguridad
    add_content_slide(
        prs,
        "Seguridad de aplicación e infraestructura",
        [
            "Aplicación: sesiones PHP, validación de entradas, separación de roles admin/usuario.",
            "Red: Security Groups en AWS; acceso restringido a Prometheus, Grafana, Alertmanager y MySQL.",
            "Simulador de carga: token de control, límites de URLs y TLS; evita abuso si se expone la UI.",
        ],
        "6:45–7:30. Resume riesgos y mitigaciones. Menciona alineación con prácticas del documento (§5.2 y simulador).",
    )

    # 11 — AWS
    add_picture_slide(
        prs,
        "Despliegue en AWS (EC2 + Docker Compose)",
        img4,
        [
            "Instancia EC2 con bootstrap (user-data) para instalar Docker y levantar el stack.",
            "Reglas de entrada mínimas para HTTP/HTTPS y UIs de monitorización según laboratorio.",
            "Documentación de costes y operación en la memoria y en docs/ del repositorio.",
        ],
        "7:30–8:30. Muestra el diagrama: un solo slide. Si hay tiempo, una captura de consola en vivo.",
        width_in=11.5,
        left_in=0.65,
    )

    # 12 — Monitorización stack
    add_content_slide(
        prs,
        "Monitorización y observabilidad",
        [
            "Prometheus recoge métricas de contenedores, MySQL, Node, blackbox y endpoint PHP (métricas app).",
            "Grafana visualiza dashboards preprovisionados (JSON en monitoring/grafana/).",
            "Alertmanager notifica incidencias (SMTP/SES según configuración).",
        ],
        "8:30–9:30. Mensaje: la observabilidad no es decorativa; alimenta operación y defensa del diseño.",
    )

    # 13 — Evidencia capturas
    add_two_picture_slide(
        prs,
        "Evidencias: Grafana y estado de contenedores",
        img19,
        img20,
        "9:30–11:00. DEMO: targets en Prometheus, dashboard Grafana, carga breve (JMeter/UI) y comenta "
        "tráfico real vs simulado, como en el plan 7.4.",
    )

    # 14 — Validación
    add_content_slide(
        prs,
        "Validación y pruebas",
        [
            "Pruebas funcionales: flujos de registro, login, citas, noticias y administración.",
            "Pruebas de despliegue: arranque reproducible con Docker y comprobación de servicios.",
            "Matriz de pruebas en memoria (§7.3): trazabilidad entre requisito y evidencia.",
        ],
        "11:00–12:00. Resume en tres frases; el detalle queda en la memoria.",
    )

    # 15 — Presupuesto
    add_content_slide(
        prs,
        "Presupuesto anual orientativo",
        [
            "Coste total orientativo documentado: ~1.146 EUR/año (infra cloud, dominio, correo, etc.).",
            "Desglose principal en la memoria (Tablas de §6.2): instancia, almacenamiento, tiempo de operación.",
            "El presupuesto acompaña decisiones técnicas con estimación explícita, no solo listado de precios.",
        ],
        "12:00–12:45. Una cifra memorable. Indica que es orientativo y depende de región/tráfico.",
    )

    # 16 — Conclusiones
    add_content_slide(
        prs,
        "Conclusiones y mejoras futuras",
        [
            "Consecución: web operativa, BD normalizada, Docker local/AWS, monitorización activa.",
            "Límites razonables del tiempo de proyecto (HTTPS definitivo, endurecimiento adicional, etc.).",
            "Mejoras propuestas en memoria: RDS, CI/CD, backups en S3, políticas RGPD.",
        ],
        "12:45–13:45. Honestidad + visión profesional. Cierra el arco objetivos→resultado→siguientes pasos.",
    )

    # 17 — Cierre
    add_content_slide(
        prs,
        "Demostración en vivo y preguntas",
        [
            "Checklist rápida: app en navegador → docker compose ps → Grafana/Prometheus → carga breve.",
            "Agradecimiento al tribunal y al tutor.",
            "¿Preguntas?",
        ],
        "Cierre (~2–4 min en guion de 20 min). Repite solo lo imprescindible de la demo si ya lo enseñaste; "
        "prioriza margen para preguntas (rúbrica: exposición + respuestas).",
    )

    # Ajustar fuentes en portada (segundo placeholder)
    for slide in prs.slides:
        set_slide_background_light(slide)

    return prs


def main():
    root = repo_root()
    out = root / "docs" / OUTPUT_NAME
    prs = build_presentation(root)
    prs.save(str(out))
    print(f"Guardado: {out}")
    missing = []
    for name in ("image1.png", "image2.png", "image3.png", "image4.png", "image19.png", "image20.png"):
        if resolve_image(root, name) is None:
            missing.append(name)
    if missing:
        print("Aviso: no se encontraron algunas imágenes:", ", ".join(missing))


if __name__ == "__main__":
    main()
